from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import os

#--------------- NOMBRAR AL ARCHIVO .pptx --------------- #
nombre_presentacion = 'presentación'
# -------------------------------------------------------------------- #

# Ruta a la carpeta que contiene las imágenes
image_folder = 'imagenes'

# Obtener lista de archivos en la carpeta con su fecha de modificación
image_files_with_dates = []
for f in os.listdir(image_folder):
    if f.lower().endswith(('png', 'jpg', 'jpeg', 'bmp', 'gif')):
        file_path = os.path.join(image_folder, f)
        modification_time = os.path.getmtime(file_path)
        # Convertir el tiempo de modificación a una fecha (sin hora)
        modification_date = datetime.fromtimestamp(modification_time).date()
        image_files_with_dates.append((f, modification_date))

# Ordenar los archivos por fecha (de más antigua a más reciente)
image_files_with_dates.sort(key=lambda x: x[1])

# Crear una lista solo con los nombres de los archivos ordenados
image_files = [f[0] for f in image_files_with_dates]
print(image_files)

prs = Presentation()

# Añadir una diapositiva por cada imagen
for image_file in image_files:
    image_path = os.path.join(image_folder, image_file)
    
    # Añadir una diapositiva en blanco
    slide_layout = prs.slide_layouts[6]  # 6 corresponde a una diapositiva vacía
    slide = prs.slides.add_slide(slide_layout)
    
    # Añadir la imagen a la diapositiva
    top = Inches(0.75)
    left = Inches(0.25)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(9.5), height=Inches(6))

# Guardar la presentación en un archivo
prs.save(f'{nombre_presentacion}.pptx')
print(f"Presentación guardada como '{nombre_presentacion}.pptx'")